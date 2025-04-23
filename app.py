# -*- coding: utf-8 -*-
"""
AI Judge - Fekrety Innovation Evaluator (No JSON output)

This Streamlit app allows you to submit two files:
  1. A Template Excel file containing project details.
  2. A Pitch Deck PPTX file containing the pitch deck information.

The app displays the submission details from the Excel file and evaluates the idea via a multi-agent system.
The complete agent thought process is shown with a typewriter effect,
and the final evaluation is displayed in Markdown and available for download as a PDF file.
"""

import os
import sys
import io
import re
import time
import base64
import pandas as pd
from dotenv import load_dotenv
import streamlit as st
import warnings
import streamlit.components.v1 as components
from reportlab.lib.pagesizes import A4
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image
from reportlab.lib.units import inch
import datetime
import zipfile
from pptx import Presentation
from lxml import etree
from io import BytesIO
import PyPDF2
from docx import Document
from crewai import Agent, Task, Crew
from langchain_openai import ChatOpenAI
from openai import OpenAI
from crewai_tools import ScrapeWebsiteTool, SerperDevTool

# --- Global Configuration ---
load_dotenv()
api_key = os.getenv("API_KEY")
os.environ["OPENAI_API_KEY"] = api_key
serper_api_key = os.getenv("SERPER_API_KEY")
os.environ["SERPER_API_KEY"] = serper_api_key
search_tool = SerperDevTool(recency_days=365)
scrape_tool = ScrapeWebsiteTool()

AVAILABLE_MODELS = {
    "GPT-4.1": "gpt-4.1",
    "GPT-4.1-mini": "gpt-4.1-mini",
    "GPT-4o-mini": "gpt-4o-mini",
    "GPT-4o": "gpt-4o",
    "GPT-o3": "o3",
    "GPT-o4-mini": "o4-mini"
}
default_model_choice = "GPT-4o"

# --- Helper Functions ---

def initialize_ai_clients(model_name):
    client = OpenAI(api_key=api_key)
    llm = ChatOpenAI(model=model_name, api_key=api_key)
    return client, llm

def clean_output(raw_text: str) -> str:
    ansi_escape = re.compile(r'\x1B\[[0-?]*[ -/]*[@-~]')
    return ansi_escape.sub('', raw_text)

def read_xlsx(file):
    df = pd.read_excel(file)
    return df


def read_pptx(file) -> str:
    """
    Extract literally *all* text from a PPTX—including SmartArt, tables, grouped shapes, etc.—by
    walking every <a:t> element in every slide part.
    """
    # load into python-pptx to get slides
    pptx_io = BytesIO(file.read())
    prs = Presentation(pptx_io)
    all_text = []

    # first, grab anything python-pptx already surfaces
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    all_text.append(para.text.strip())

    # now, as a catch-all, scan *all* the slide XMLs for <a:t> nodes
    pptx_io.seek(0)
    with zipfile.ZipFile(pptx_io) as z:
        # any XML under ppt/slides or ppt/diagrams
        for name in z.namelist():
            if name.endswith(".xml") and (name.startswith("ppt/slides") or name.startswith("ppt/diagrams")):
                xml = z.read(name)
                root = etree.fromstring(xml)
                # the DrawingML namespace for text is always this
                for t in root.xpath("//a:t", namespaces={"a":"http://schemas.openxmlformats.org/drawingml/2006/main"}):
                    if t.text and t.text.strip():
                        all_text.append(t.text.strip())

    # de‐duplicate and join
    # (you can skip dedupe if order matters)
    seen = set()
    cleaned = []
    for line in all_text:
        if line not in seen:
            seen.add(line)
            cleaned.append(line)

    return "\n".join(cleaned)

def read_pdf(file):
    pdf_reader = PyPDF2.PdfReader(file)
    text = ""
    for page in pdf_reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text
    return text

def read_docx(file):
    doc = Document(file)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def read_txt(file):
    return file.read().decode("utf-8")

def create_pdf_report(idea_title, score, justification, template_df):
    """
    Create a PDF report with evaluation results - with four separate tables:
    - First table with 3 columns
    - Three subsequent tables with 2 columns each
    All arranged vertically one below another
    """
    buffer = io.BytesIO()
    # Set page margins to ensure content fits properly
    doc = SimpleDocTemplate(buffer, pagesize=A4, 
                           rightMargin=72, leftMargin=72,
                           topMargin=72, bottomMargin=72)
    styles = getSampleStyleSheet()
    
    # Create a custom style for the title
    title_style = ParagraphStyle(
        'Title',
        parent=styles['Heading1'],
        fontSize=18,
        alignment=1,  # Center alignment
        spaceAfter=12
    )
    
    # Create a custom style for section headings
    section_style = ParagraphStyle(
        'Section',
        parent=styles['Heading2'],
        fontSize=14,
        spaceBefore=10,
        spaceAfter=6
    )
    
    # Create a custom style for the score display
    score_style = ParagraphStyle(
        'Score',
        parent=styles['Normal'],
        fontSize=24,
        alignment=1,  # Center alignment
        textColor=colors.blue
    )
    
    # Build the PDF content
    elements = []
    
    # Add logo if available
    try:
        logo = Image("judge_bg.png", width=1.5*inch, height=1.5*inch)
        elements.append(logo)
    except:
        pass  # Continue without logo if not found
    
    # Add title and date
    elements.append(Paragraph("AI Judge - Innovation Evaluation Report", title_style))
    elements.append(Paragraph(f"Generated on: {datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", styles["Normal"]))
    elements.append(Spacer(1, 0.25*inch))
    
    # Add idea title
    elements.append(Paragraph(f"Idea: {idea_title}", section_style))
    elements.append(Spacer(1, 0.1*inch))
    
    # Add score
    elements.append(Paragraph(f"Overall Score: {score}/100", score_style))
    elements.append(Spacer(1, 0.25*inch))
    
    # Add submission details section heading
    elements.append(Paragraph("Submission Details:", section_style))
    
    # Get all columns from the dataframe
    all_columns = list(template_df.columns)
    
    # Common table style
    table_style = TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.lightblue),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, -1), 10),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 8),
        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
        ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ('WORDWRAP', (0, 0), (-1, -1), True),
        ('LEFTPADDING', (0, 0), (-1, -1), 6),
        ('RIGHTPADDING', (0, 0), (-1, -1), 6)
    ])
    
    total_cols = len(all_columns)
    cols_used = 0
    
    # Table 1: First 3 columns (or fewer if not available)
    cols_table1 = min(3, total_cols)
    table1_cols = all_columns[:cols_table1]
    cols_used += cols_table1
    
    # Table 2: Next 2 columns (or fewer if not available)
    cols_table2 = min(2, max(0, total_cols - cols_used))
    table2_cols = all_columns[cols_used:cols_used+cols_table2] if cols_table2 > 0 else []
    cols_used += cols_table2
    
    # Table 3: Next 2 columns (or fewer if not available)
    cols_table3 = min(2, max(0, total_cols - cols_used))
    table3_cols = all_columns[cols_used:cols_used+cols_table3] if cols_table3 > 0 else []
    cols_used += cols_table3
    
    # Table 4: Next 2 columns (or fewer if not available)
    cols_table4 = min(2, max(0, total_cols - cols_used))
    table4_cols = all_columns[cols_used:cols_used+cols_table4] if cols_table4 > 0 else []
    
    def create_table(table_cols):
        if not table_cols:
            return None
        table_data = [table_cols]
        for _, row in template_df.iterrows():
            data_row = []
            for col in table_cols:
                data_row.append(str(row[col]))
            table_data.append(data_row)
        col_width = doc.width / len(table_cols)
        table = Table(table_data, colWidths=[col_width] * len(table_cols))
        table.setStyle(table_style)
        return table
    
    tables = [
        (table1_cols, "Table 1 (Primary Details)"),
        (table2_cols, "Table 2 (Additional Details)"),
        (table3_cols, "Table 3 (Additional Details)"),
        (table4_cols, "Table 4 (Additional Details)")
    ]
    
    tables_added = 0
    for i, (cols, table_name) in enumerate(tables):
        table = create_table(cols)
        if table:
            tables_added += 1
            elements.append(table)
            elements.append(Spacer(1, 0.15*inch))
    
    if tables_added == 0:
        elements.append(Paragraph("No submission details available", styles["Normal"]))
    
    elements.append(Spacer(1, 0.25*inch))
    elements.append(Paragraph("Evaluation Justification:", section_style))
    
    justification_style = ParagraphStyle(
        'Justification',
        parent=styles['Normal'],
        fontSize=10,
        spaceAfter=6,
        leading=14
    )
    
    formatted_sections = []
    pattern = r"(Creativity|Viability|NSV):?\s*(\d+/\d+)?\s*-?\s*(.*?)(?=(Creativity|Viability|NSV):|$)"
    matches = re.findall(pattern, justification, re.DOTALL)
    
    if matches:
        for section, score_part, content, _ in matches:
            score_text = f"{score_part} - " if score_part else ""
            section_text = f"<b>{section} (Score Component):</b><br/>{score_text}{content.strip()}"
            formatted_sections.append(Paragraph(section_text, justification_style))
    else:
        sections = {"Creativity": "", "Viability": "", "NSV": ""}
        current_section = None
        for line in justification.split(", "):
            for section in sections.keys():
                if line.startswith(section + ":"):
                    current_section = section
                    content = line[len(section) + 1:].strip()
                    sections[section] = content
                    break
            if current_section and not line.startswith(tuple([s + ":" for s in sections.keys()])):
                sections[current_section] += ", " + line
        for section, content in sections.items():
            if content:
                section_text = f"<b>{section} (Score Component):</b><br/>{content}"
                formatted_sections.append(Paragraph(section_text, justification_style))
    
    for section_paragraph in formatted_sections:
        elements.append(section_paragraph)
        elements.append(Spacer(1, 0.1*inch))
    
    elements.append(Spacer(1, 0.5*inch))
    elements.append(Paragraph("This evaluation was generated by AI Judge - Gulf Bank's Innovation Evaluator", ParagraphStyle("Footer", parent=styles["Italic"], alignment=1)))
    
    doc.build(elements)
    buffer.seek(0)
    return buffer

def process_idea(idea_text, llm, excel_title):
    # Define agents
    summary_agent = Agent(
        role="Idea Summarizer",
        goal="Summarize the core concept and objectives of the submitted idea.",
        verbose=True,
        allow_delegation=False,
        backstory="You are an executive summary writer with banking expertise in Kuwait."
    )
    analysis_agent = Agent(
        role="Idea Analyzer",
        goal=(
            "Analyze Creativity, Viability, and NSV. "
            "Use internet tools for context if needed, but present insights clearly."
        ),
        verbose=True,
        allow_delegation=False,
        backstory="You are a Gulf Bank expert evaluating innovations.",
        tools=[scrape_tool, search_tool]
    )
    scorer_agent = Agent(
        role="Idea Scorer",
        goal=(
            "Score the idea on a scale of 1 to 100 and provide a detailed breakdown. "
            "Output exactly in this format (no JSON):\n"
            "Score: <number>/100\n"
            "Creativity: <explanation>\n"
            "Viability: <explanation>\n"
            "NSV: <explanation>"
        ),
        verbose=True,
        allow_delegation=False,
        backstory="You are a tough innovation judge for Gulf Bank."
    )

    # Tasks
    task1 = Task(
        description="Generate a summary of the idea.",
        expected_output="A concise summary paragraph.",
        agent=summary_agent
    )
    task2 = Task(
        description="Analyze creativity, viability, and NSV based on the summary.",
        expected_output="Detailed analysis notes.",
        agent=analysis_agent,
        context=[task1]
    )
    task3 = Task(
        description="Score the idea and justify each criterion.",
        expected_output="Plain-text formatted score and breakdown.",
        agent=scorer_agent,
        context=[task2]
    )

    crew = Crew(
        agents=[summary_agent, analysis_agent, scorer_agent],
        tasks=[task1, task2, task3],
        manager=llm,
        verbose=True
    )

    # Run
    buf = io.StringIO()
    old_stdout = sys.stdout
    sys.stdout = buf
    crew.kickoff(inputs={"idea_text": idea_text})
    sys.stdout = old_stdout
    thought = clean_output(buf.getvalue())
    buf.close()

    # Extract score and justification
    # Expect lines: Score: ##/100, then Creativity:, Viability:, NSV:
    lines = thought.strip().splitlines()
    score_line = next((l for l in lines if l.startswith("Score:")), None)
    score = None
    if score_line:
        m = re.search(r"Score:\s*(\d+)/(\d+)", score_line)
        if m:
            score = int(m.group(1))
    # Collect justification lines
    just_lines = [l for l in lines if l.startswith(("Creativity:", "Viability:", "NSV:"))]
    justification = "\n".join(just_lines)

    return excel_title, score, justification, thought

# --- Streamlit UI Setup ---
st.set_page_config(
    page_title="AI Judge - Innovation Evaluator", 
    page_icon="⚖️",
    layout="wide",
    initial_sidebar_state="expanded"
)

st.markdown("""
<style>
    .stApp {
        background-color: #f8f9fa;
    }
    h1, h2, h3, h4, h5, h6 {
        color: #2c3e50 !important;
        font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
    }
    .stButton>button {
        background-color: #4a6fa5;
        color: white;
        border-radius: 8px;
        padding: 0.5rem 1rem;
        border: none;
        font-weight: 500;
        transition: all 0.3s ease;
    }
    .stButton>button:hover {
        background-color: #3a5a8f;
        transform: translateY(-2px);
        box-shadow: 0 4px 8px rgba(0,0,0,0.1);
    }
    .stFileUploader>div>div {
        border: 2px dashed #4a6fa5;
        border-radius: 12px;
        padding: 2rem;
        background-color: rgba(74, 111, 165, 0.05);
    }
    .stExpander {
        border: 1px solid #e1e4e8;
        border-radius: 12px;
        box-shadow: 0 2px 8px rgba(0,0,0,0.1);
    }
    .stDataFrame {
        border-radius: 12px;
        box-shadow: 0 2px 8px rgba(0,0,0,0.1);
    }
    .stSuccess {
        border-radius: 12px;
    }
    .custom-card {
        background-color: white;
        border-radius: 12px;
        padding: 1.5rem;
        box-shadow: 0 4px 6px rgba(0,0,0,0.1);
        margin-bottom: 1.5rem;
        border-left: 4px solid #4a6fa5;
    }
    .agent-thought {
        background-color: #f0f4f8;
        border-radius: 12px;
        padding: 1rem;
        margin: 0.5rem 0;
        font-family: 'Courier New', monospace;
        font-size: 14px;
        line-height: 1.6;
    }
    .score-display {
        font-size: 2.5rem;
        font-weight: bold;
        color: #4a6fa5;
        text-align: center;
        margin: 1rem 0;
    }
</style>
""", unsafe_allow_html=True)

col1, col2 = st.columns([1, 3])
with col1:
    st.image("judge_bg.png", width=150)
with col2:
    st.title("⚖️ AI Judge - Innovation Evaluator")
    st.caption("Your AI-powered innovation assessment tool")

with st.container():
    st.markdown("""
    <div class="custom-card">
        <h3 style="color: #2c3e50; margin-top: 0;">Welcome to the Future of Innovation Evaluation! 🚀</h3>
        <p>Our cutting-edge AI delivers an unbiased, detailed evaluation based on:</p>
        <ul>
            <li><b>Creativity</b> (35 points): Originality and uniqueness</li>
            <li><b>Viability</b> (30 points): Practical feasibility</li>
            <li><b>NSV</b> (35 points): Need, Solution, and Value</li>
        </ul>
    </div>
    """, unsafe_allow_html=True)

with st.container():
    st.markdown("""
    <div class="custom-card">
        <h4 style="color: #2c3e50; margin-top: 0;">Evaluation Settings</h4>
    """, unsafe_allow_html=True)
    model_choice = st.selectbox("Select AI Model", list(AVAILABLE_MODELS.keys()),
                                index=list(AVAILABLE_MODELS.keys()).index(default_model_choice),
                                help="Choose the AI model that will evaluate your submission")
    st.markdown("</div>", unsafe_allow_html=True)

col1, col2 = st.columns(2)
with col1:
    with st.container():
        st.markdown("""
        <div class="custom-card">
            <h4 style="color: #2c3e50; margin-top: 0;">Project Template</h4>
        """, unsafe_allow_html=True)
        template_file = st.file_uploader("Upload Excel Template", type=["xlsx"],
                                         label_visibility="collapsed")
        st.markdown("</div>", unsafe_allow_html=True)

with col2:
    with st.container():
        st.markdown("""
        <div class="custom-card">
            <h4 style="color: #2c3e50; margin-top: 0;">Pitch Deck</h4>
        """, unsafe_allow_html=True)
        pitch_file = st.file_uploader("Upload PowerPoint Pitch", type=["pptx"],
                                      label_visibility="collapsed")
        st.markdown("</div>", unsafe_allow_html=True)

if template_file is not None and pitch_file is not None:
    with st.container():
        st.markdown("""<div class="custom-card">
            <h4 style="color: #2c3e50; margin-top: 0;">Submission Details</h4>""", unsafe_allow_html=True)
        template_df = pd.read_excel(template_file)
        st.dataframe(template_df)
        st.markdown("</div>", unsafe_allow_html=True)
        
        pitch_text = read_pptx(pitch_file)
        st.markdown("### Pitch Deck Details (Extracted Text)")
        st.text_area("Pitch Deck Content", pitch_text, height=300)
    
    if st.button("Evaluate Idea", type="primary"):
        with st.spinner("Evaluating your idea..."):
            client, llm = initialize_ai_clients(AVAILABLE_MODELS[model_choice])
            # Extract the idea title from the Excel file using the "Project Title" header
            excel_title = (template_df["Project Title"].iloc[0] 
                           if "Project Title" in template_df.columns and not template_df["Project Title"].empty 
                           else "Untitled")
            idea_title, score, justification, thought_process = process_idea(pitch_text, llm, excel_title)
            
            st.markdown("### AI Deliberation Process")
            with st.expander("🔍 View Detailed Agent Thought Process", expanded=True):
                with st.chat_message("assistant"):
                    text_placeholder = st.empty()
                    displayed_text = ""
                    for char in thought_process:
                        displayed_text += char
                        processed_text = displayed_text.replace(
                            "Agent:", "<span style='font-weight:bold;color:#2c3e50'>Agent:</span>"
                        ).replace(
                            "Thought:", "<span style='font-weight:bold;color:#2c3e50'>Thought:</span>"
                        ).replace(
                            "Final Answer:", "<span style='font-weight:bold;color:#2c3e50'>Final Answer:</span>"
                        )
                        text_placeholder.markdown(f"""
                        <div style="font-family: 'Courier New', monospace; font-size: 14px; background-color: #f0f4f8; padding: 1rem; border-radius: 12px; margin: 0.5rem 0; line-height: 1.6;">
                            {processed_text}▌
                        </div>
                        """, unsafe_allow_html=True)
                        time.sleep(0.002)
                    
                    text_placeholder.markdown(f"""
                    <div style="font-family: 'Courier New', monospace; font-size: 14px; background-color: #f0f4f8; padding: 1rem; border-radius: 12px; margin: 0.5rem 0; line-height: 1.6;">
                        {processed_text}
                    </div>
                    """, unsafe_allow_html=True)
            
            st.success("Evaluation Complete!")
            formatted_justification = justification.replace("Creativity:", "\n\n**Creativity:**") \
                                                 .replace("Viability:", "\n\n**Viability:**") \
                                                 .replace("NSV:", "\n\n**NSV:**")
            
            st.markdown(f"""
            <div class="custom-card">
                <h2 style="color: #2c3e50; text-align: center; margin-bottom: 0.5rem;">Evaluation Results</h2>
                <h3 style="text-align: center; margin-top: 0;">{idea_title}</h3>
                <div class="score-display" style="margin: 1rem 0;">{score}/100</div>
                <p style="font-weight: bold; color: #2c3e50; margin-bottom: 0.5rem;">Detailed Justification:</p>
                <div style="line-height: 1.8;">
                    {formatted_justification}
                </div>
            </div>
            """, unsafe_allow_html=True)
            
            pdf_buffer = create_pdf_report(idea_title, score, justification, template_df)
            
            st.download_button(
                label="📥 Download Full Report (PDF)",
                data=pdf_buffer,
                file_name=f"evaluation_{idea_title.replace(' ','_')}.pdf",
                mime="application/pdf"
            )
            

            
